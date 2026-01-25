"""
title: PPT/PDF Vision Filter
author: GLChemTec
version: 11.0
description: Processes PPT/PPTX/PDF files for vision analysis. Only processes files in the CURRENT message - ignores conversation history.
"""

import os
import base64
import tempfile
import shutil
import subprocess
import zipfile
import hashlib
import time
from typing import Optional, List, Dict, Any, Set

from pydantic import BaseModel, Field

# Optional: local PPTX text/table extraction
try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except Exception:
    PPTX_AVAILABLE = False


class Filter:
    class Valves(BaseModel):
        priority: int = Field(default=0, description="Filter priority (0 = highest)")
        enabled: bool = Field(default=True, description="Enable PPT/PDF vision processing")
        debug: bool = Field(default=True, description="Enable debug logging")

        # Rendering quality (high DPI for spectra clarity)
        dpi: int = Field(default=600, description="DPI for PDF rendering (600 = high quality for spectra/NMR)")
        max_pages: int = Field(default=30, description="Max pages/slides to render")
        output_format: str = Field(default="png", description="png or jpeg")
        jpeg_quality: int = Field(default=85, description="JPEG quality if using jpeg")

        # Size limits
        max_total_image_mb: float = Field(default=50.0, description="Max total image payload (MB)")
        
        # Timeouts (balanced for speed and completeness)
        libreoffice_base_timeout: int = Field(default=30, description="Base LibreOffice timeout (sec)")
        libreoffice_per_slide_timeout: int = Field(default=3, description="Additional seconds per slide")
        max_timeout: int = Field(default=120, description="Maximum timeout cap (120s) - allows more time for large PPTs")
        max_processing_time: int = Field(default=120, description="Max total processing time (sec) - allows time for 10-25 slide PPTs")

        # Features
        extract_text: bool = Field(default=True, description="Extract text/tables from PPTX")
        extract_embedded_images: bool = Field(default=True, description="Extract PNG/JPG images from PPTX")
        convert_emf_wmf: bool = Field(default=True, description="Attempt EMF/WMF conversion")

    def __init__(self):
        # CRITICAL: Always ensure valves exists, even if init fails
        # This prevents OpenWebUI from crashing when introspecting filters
        try:
            self.valves = self.Valves()
            # Track processed files to avoid re-processing in same session
            self._processed_files: Set[str] = set()
        except Exception as e:
            # If initialization fails, disable the filter to prevent crashes
            print(f"[PPT-PDF-VISION] ERROR in __init__: {e}")
            import traceback
            print(f"[PPT-PDF-VISION] Traceback: {traceback.format_exc()}")
            # Create a minimal disabled filter - MUST succeed or OpenWebUI crashes
            try:
                self.valves = self.Valves(enabled=False)
                self._processed_files: Set[str] = set()
            except Exception as e2:
                # Last resort - this should never happen
                print(f"[PPT-PDF-VISION] CRITICAL: Cannot create Valves - {e2}")
                raise

    def _log(self, msg: str) -> None:
        if self.valves.debug:
            print(f"[PPT-PDF-VISION] {msg}")

    # =========================================================================
    # FILE DISCOVERY - Only current message, not history
    # =========================================================================
    
    def _get_file_path(self, file_obj: Dict[str, Any]) -> str:
        """Extract file path from various OpenWebUI file object formats."""
        if not isinstance(file_obj, dict):
            return ""
        
        # Nested file object
        if isinstance(file_obj.get("file"), dict):
            f = file_obj["file"]
            path = (f.get("path") or "").strip()
            if path:
                return path
            if isinstance(f.get("meta"), dict):
                return (f["meta"].get("path") or "").strip()
        
        # Direct path
        return (file_obj.get("path") or "").strip()

    def _get_file_name(self, file_obj: Dict[str, Any]) -> str:
        """Extract filename from file object."""
        if not isinstance(file_obj, dict):
            return ""
        
        if isinstance(file_obj.get("file"), dict):
            f = file_obj["file"]
            name = f.get("filename") or f.get("name") or ""
            if not name and isinstance(f.get("meta"), dict):
                name = f["meta"].get("name") or ""
            return (name or "").lower().strip()
        
        return (file_obj.get("name") or file_obj.get("filename") or "").lower().strip()

    def _get_file_id(self, file_obj: Dict[str, Any]) -> str:
        """Get unique file ID for tracking."""
        if isinstance(file_obj.get("file"), dict):
            f = file_obj["file"]
            return f.get("id") or f.get("path") or ""
        return file_obj.get("id") or file_obj.get("path") or ""

    def _get_current_message_files(self, body: dict, messages: list) -> List[Dict[str, Any]]:
        """
        Get files from CURRENT message only.
        This is the key to not re-processing old files.
        """
        files: List[Dict[str, Any]] = []
        
        # Files in request body (current upload)
        if isinstance(body.get("files"), list):
            files.extend(body["files"])
        
        # Files in the last message (if it's a user message with attachments)
        if messages:
            last_msg = messages[-1]
            if last_msg.get("role") == "user":
                if isinstance(last_msg.get("files"), list):
                    files.extend(last_msg["files"])
                if isinstance(last_msg.get("attachments"), list):
                    files.extend(last_msg["attachments"])
        
        # Deduplicate by path
        seen = set()
        unique = []
        for f in files:
            path = self._get_file_path(f)
            if path and path not in seen:
                seen.add(path)
                unique.append(f)
        
        return unique

    def _is_supported_file(self, file_name: str) -> bool:
        """Check if file type is supported by this filter."""
        return file_name.endswith((".ppt", ".pptx", ".pdf"))

    def _file_hash(self, file_path: str) -> str:
        """Generate hash of file for dedup tracking."""
        try:
            with open(file_path, "rb") as f:
                return hashlib.md5(f.read(8192)).hexdigest()  # First 8KB
        except:
            return file_path

    # =========================================================================
    # LIBREOFFICE CONVERSION
    # =========================================================================

    def _find_libreoffice(self) -> Optional[str]:
        """Find LibreOffice or unoconv for conversion."""
        # Try unoconv first (faster, uses persistent listener)
        unoconv = shutil.which("unoconv")
        if unoconv:
            return unoconv
        # Fallback to direct LibreOffice
        return shutil.which("libreoffice") or shutil.which("soffice")

    def _count_slides(self, ppt_path: str) -> int:
        """Count slides for timeout calculation."""
        try:
            if PPTX_AVAILABLE:
                prs = Presentation(ppt_path)
                return len(prs.slides)
            else:
                with zipfile.ZipFile(ppt_path, 'r') as z:
                    slides = [f for f in z.namelist() 
                             if f.startswith('ppt/slides/slide') and f.endswith('.xml')]
                    return len(slides)
        except:
            return 10

    def _convert_pptx_to_pdf(self, ppt_path: str, out_dir: str) -> Optional[str]:
        """Convert PPTX to PDF using LibreOffice with optimizations and diagnostics."""
        import time
        start_time = time.time()
        
        lo = self._find_libreoffice()
        if not lo:
            self._log("LibreOffice not found - skipping PDF conversion")
            return None

        # Calculate timeout based on slide count
        slide_count = self._count_slides(ppt_path)
        timeout = min(
            self.valves.libreoffice_base_timeout + (slide_count * self.valves.libreoffice_per_slide_timeout),
            self.valves.max_timeout
        )
        
        self._log(f"Converting PPTX ({slide_count} slides) with {timeout}s timeout")
        self._log(f"File size: {os.path.getsize(ppt_path) / 1024 / 1024:.1f}MB")

        profile_dir = tempfile.mkdtemp(prefix="lo_profile_")
        try:
            env = os.environ.copy()
            env["HOME"] = "/tmp"
            env["TMPDIR"] = "/tmp"
            # Optimize LibreOffice performance
            env["SAL_USE_VCLPLUGIN"] = "headless"  # Force headless rendering
            env["SAL_DISABLE_OPENCL"] = "1"  # Disable OpenCL (can cause issues)
            env["SAL_DISABLE_OPENGL"] = "1"  # Disable OpenGL (can cause issues)

            # Use optimized LibreOffice flags with PDF export parameters for 600 DPI
            # Format: pdf:impress_pdf_Export:{"MaxImageResolution":600,"ReduceImageResolution":false}
            pdf_params = (
                'pdf:impress_pdf_Export:'
                '{"MaxImageResolution":{"type":"long","value":"600"},'
                '"ReduceImageResolution":{"type":"boolean","value":"false"},'
                '"UseLosslessCompression":{"type":"boolean","value":"true"}}'
            )
            
            cmd = [
                lo,
                "--headless",
                "--nologo",
                "--nofirststartwizard",
                "--nodefault",  # Don't load default documents
                "--nolockcheck",  # Skip lock file checking (faster)
                f"-env:UserInstallation=file://{profile_dir}",
                "--convert-to", pdf_params,
                "--outdir", out_dir,
                ppt_path,
            ]
            
            self._log(f"Starting LibreOffice conversion...")
            conversion_start = time.time()
            
            result = subprocess.run(cmd, timeout=timeout, capture_output=True, text=True, env=env)
            
            conversion_time = time.time() - conversion_start
            self._log(f"LibreOffice conversion took {conversion_time:.1f}s (returncode: {result.returncode})")
            
            if result.returncode != 0:
                self._log(f"LibreOffice error (code {result.returncode})")
                if result.stderr:
                    self._log(f"stderr: {result.stderr[:500]}")
                if result.stdout:
                    self._log(f"stdout: {result.stdout[:500]}")

            # Find output PDF
            base = os.path.splitext(os.path.basename(ppt_path))[0]
            expected = os.path.join(out_dir, base + ".pdf")
            if os.path.exists(expected):
                return expected

            for fn in os.listdir(out_dir):
                if fn.lower().endswith(".pdf"):
                    return os.path.join(out_dir, fn)

            return None

        except subprocess.TimeoutExpired:
            self._log(f"LibreOffice timed out after {timeout}s - killing process")
            try:
                subprocess.run(["pkill", "-9", "-f", "soffice"], timeout=2, capture_output=True)
            except:
                pass
            return None
        except Exception as e:
            self._log(f"Conversion error: {e}")
            return None
        finally:
            shutil.rmtree(profile_dir, ignore_errors=True)

    # =========================================================================
    # PDF TO IMAGES
    # =========================================================================

    def _convert_pdf_to_images(self, pdf_path: str, output_dir: str) -> List[str]:
        """Render PDF pages to images with timing diagnostics."""
        import time
        start_time = time.time()
        paths = []
        try:
            from pdf2image import convert_from_path

            fmt = self.valves.output_format.lower()
            if fmt not in ("png", "jpeg", "jpg"):
                fmt = "png"

            pdf_size = os.path.getsize(pdf_path) / 1024 / 1024
            self._log(f"Rendering PDF at {self.valves.dpi} DPI (PDF size: {pdf_size:.1f}MB)")

            render_start = time.time()
            images = convert_from_path(
                pdf_path,
                dpi=self.valves.dpi,
                first_page=1,
                last_page=self.valves.max_pages,
                thread_count=1,  # Single thread to avoid memory issues
            )
            render_time = time.time() - render_start
            self._log(f"PDF2image conversion took {render_time:.1f}s for {len(images)} pages")

            max_bytes = int(self.valves.max_total_image_mb * 1024 * 1024)
            total_bytes = 0

            save_start = time.time()
            for i, img in enumerate(images, start=1):
                if fmt == "png":
                    out_path = os.path.join(output_dir, f"page_{i:03d}.png")
                    img.save(out_path, format="PNG")
                else:
                    out_path = os.path.join(output_dir, f"page_{i:03d}.jpg")
                    img.convert("RGB").save(out_path, format="JPEG", quality=self.valves.jpeg_quality)

                size = os.path.getsize(out_path)
                total_bytes += size
                paths.append(out_path)

                if total_bytes > max_bytes:
                    self._log(f"Size limit reached at page {i}")
                    break

            save_time = time.time() - save_start
            total_time = time.time() - start_time
            self._log(f"Rendered {len(paths)} pages ({total_bytes/1024/1024:.1f}MB)")
            self._log(f"Timing: conversion={render_time:.1f}s, save={save_time:.1f}s, total={total_time:.1f}s")

        except Exception as e:
            self._log(f"PDF rendering error: {e}")

        return paths

    # =========================================================================
    # PPTX CONTENT EXTRACTION
    # =========================================================================

    def _extract_pptx_text(self, ppt_path: str) -> str:
        """Extract text and tables from PPTX."""
        if not PPTX_AVAILABLE or not self.valves.extract_text:
            return ""

        try:
            prs = Presentation(ppt_path)
            parts = []

            for idx, slide in enumerate(prs.slides, start=1):
                slide_text = []
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text = shape.text.strip()
                        if text:
                            slide_text.append(text)

                    if getattr(shape, "has_table", False):
                        try:
                            rows = []
                            for row in shape.table.rows:
                                cells = [cell.text.strip() for cell in row.cells]
                                rows.append(" | ".join(cells))
                            if rows:
                                slide_text.append("[Table]\n" + "\n".join(rows))
                        except:
                            pass

                if slide_text:
                    parts.append(f"--- Slide {idx} ---\n" + "\n".join(slide_text))

            return "\n\n".join(parts)

        except Exception as e:
            self._log(f"Text extraction error: {e}")
            return ""

    def _extract_pptx_images(self, ppt_path: str, output_dir: str) -> List[str]:
        """Extract embedded images (PNG, JPG) from PPTX."""
        if not self.valves.extract_embedded_images:
            return []

        paths = []
        try:
            with zipfile.ZipFile(ppt_path, 'r') as z:
                media_files = [f for f in z.namelist() 
                              if f.startswith('ppt/media/') and 
                              f.lower().endswith(('.png', '.jpg', '.jpeg', '.gif'))]

                for i, media_file in enumerate(media_files):
                    try:
                        ext = os.path.splitext(media_file)[1].lower()
                        out_path = os.path.join(output_dir, f"image_{i}{ext}")
                        
                        data = z.read(media_file)
                        
                        # Skip tiny images (likely icons)
                        if len(data) < 5000:
                            continue
                            
                        with open(out_path, 'wb') as f:
                            f.write(data)
                        
                        paths.append(out_path)
                        self._log(f"Extracted {os.path.basename(media_file)}")
                        
                    except Exception as e:
                        self._log(f"Error extracting {media_file}: {e}")

        except Exception as e:
            self._log(f"Image extraction error: {e}")

        return paths

    def _extract_emf_wmf(self, ppt_path: str, output_dir: str) -> List[str]:
        """Extract and convert EMF/WMF images from PPTX."""
        if not self.valves.convert_emf_wmf:
            return []

        paths = []
        lo = self._find_libreoffice()
        
        try:
            with zipfile.ZipFile(ppt_path, 'r') as z:
                emf_files = [f for f in z.namelist() 
                            if f.startswith('ppt/media/') and 
                            f.lower().endswith(('.emf', '.wmf'))]

                if not emf_files:
                    return []

                self._log(f"Found {len(emf_files)} EMF/WMF files")

                for i, emf_file in enumerate(emf_files):
                    try:
                        ext = os.path.splitext(emf_file)[1].lower()
                        temp_emf = os.path.join(output_dir, f"temp_{i}{ext}")
                        output_png = os.path.join(output_dir, f"emf_{i}.png")

                        # Extract EMF file
                        with z.open(emf_file) as src:
                            with open(temp_emf, 'wb') as dst:
                                dst.write(src.read())

                        # Try ImageMagick first (faster)
                        # Use lower DPI for EMF (chemical structures are line drawings, don't need 600 DPI)
                        emf_dpi = 150  # Much faster for line drawings
                        convert_cmd = shutil.which("convert") or shutil.which("magick")
                        if convert_cmd:
                            result = subprocess.run(
                                [convert_cmd, "-density", str(emf_dpi), temp_emf, 
                                 "-background", "white", "-flatten", output_png],
                                capture_output=True, timeout=5  # Reduced to 5s - fail fast
                            )
                            if result.returncode == 0 and os.path.exists(output_png):
                                paths.append(output_png)
                                self._log(f"Converted {os.path.basename(emf_file)} via ImageMagick")
                                continue

                        # Try LibreOffice: EMF -> PDF -> PNG
                        if lo:
                            pdf_dir = os.path.join(output_dir, f"emf_pdf_{i}")
                            os.makedirs(pdf_dir, exist_ok=True)
                            profile = tempfile.mkdtemp(prefix="lo_emf_")
                            
                            try:
                                env = os.environ.copy()
                                env["HOME"] = "/tmp"
                                
                                subprocess.run(
                                    [lo, "--headless", "--nologo",
                                     f"-env:UserInstallation=file://{profile}",
                                     "--convert-to", "pdf", "--outdir", pdf_dir, temp_emf],
                                    capture_output=True, timeout=5, env=env  # Reduced to 5s - fail fast
                                )

                                # Find PDF and convert to PNG
                                for fn in os.listdir(pdf_dir):
                                    if fn.endswith(".pdf"):
                                        pdf_path = os.path.join(pdf_dir, fn)
                                        try:
                                            from pdf2image import convert_from_path
                                            # Use lower DPI for EMF (chemical structures are line drawings, don't need 600 DPI)
                                            emf_dpi = 150  # Much faster for line drawings
                                            imgs = convert_from_path(pdf_path, dpi=emf_dpi, first_page=1, last_page=1)
                                            if imgs:
                                                imgs[0].save(output_png, format="PNG")
                                                if os.path.exists(output_png):
                                                    paths.append(output_png)
                                                    self._log(f"Converted {os.path.basename(emf_file)} via LibreOffice")
                                        except:
                                            pass
                                        break
                            finally:
                                shutil.rmtree(profile, ignore_errors=True)
                                shutil.rmtree(pdf_dir, ignore_errors=True)

                    except Exception as e:
                        self._log(f"EMF conversion error: {e}")

        except Exception as e:
            self._log(f"EMF extraction error: {e}")

        return paths

    # =========================================================================
    # IMAGE UTILITIES
    # =========================================================================

    def _to_data_url(self, img_path: str, max_size_mb: float = 2.0) -> Optional[str]:
        """Convert image file to base64 data URL. Faster with size limits."""
        try:
            if not os.path.exists(img_path):
                return None
            
            # Check file size - skip if too large (reduced to 2MB for speed)
            size_mb = os.path.getsize(img_path) / (1024 * 1024)
            if size_mb > max_size_mb:
                self._log(f"Skipping large image: {size_mb:.1f}MB (limit: {max_size_mb}MB)")
                return None
            
            with open(img_path, "rb") as f:
                data = f.read()
            
            # Skip if still too large after reading (safety check)
            if len(data) > max_size_mb * 1024 * 1024:
                self._log(f"Skipping image after read: {len(data)/(1024*1024):.1f}MB")
                return None
            
            ext = os.path.splitext(img_path)[1].lower()
            mime = {
                ".png": "image/png",
                ".jpg": "image/jpeg",
                ".jpeg": "image/jpeg",
                ".gif": "image/gif",
                ".webp": "image/webp",
            }.get(ext, "image/png")
            
            b64 = base64.b64encode(data).decode("utf-8")
            return f"data:{mime};base64,{b64}"
        except Exception as e:
            self._log(f"Error encoding image: {e}")
            return None

    # =========================================================================
    # MAIN FILTER
    # =========================================================================

    def inlet(self, body: dict, __user__: Optional[dict] = None) -> dict:
        """Process PPT/PDF files in the current message."""
        
        if not self.valves.enabled:
            return body

        messages = body.get("messages", [])
        if not messages:
            return body
        
        self._log("=" * 60)
        self._log("INLET - PPT/PDF Vision Filter v11.0 (PPT->PDF->PNG + EMF/WMF)")

        # Get files from CURRENT message only
        files = self._get_current_message_files(body, messages)
        
        # Filter to supported file types
        supported_files = []
        for f in files:
            name = self._get_file_name(f)
            path = self._get_file_path(f)
            
            if not path or not os.path.exists(path):
                continue
                
            if not self._is_supported_file(name):
                self._log(f"Skipping unsupported file: {name}")
                continue
            
            # Check if already processed (by file hash)
            file_hash = self._file_hash(path)
            if file_hash in self._processed_files:
                self._log(f"Skipping already processed: {name}")
                continue
            
            supported_files.append((f, name, path, file_hash))

        if not supported_files:
            self._log("No new PPT/PDF files to process")
            return body

        self._log("=" * 60)
        self._log(f"Processing {len(supported_files)} file(s)")

        all_text = []
        all_images = []

        for file_obj, file_name, file_path, file_hash in supported_files:
            self._log(f"Processing: {file_name}")
            
            with tempfile.TemporaryDirectory() as tmp:
                is_pptx = file_name.endswith((".ppt", ".pptx"))
                is_pdf = file_name.endswith(".pdf")

                # PPTX processing - FAST PATH: extract text/images first, then try PDF
                if is_pptx:
                    import time
                    start_time = time.time()
                    
                    # Count slides for timeout calculation
                    slide_count = self._count_slides(file_path)
                    self._log(f"Processing PPTX: {slide_count} slides")
                    
                    # Extract text (FAST - always do this first)
                    text = self._extract_pptx_text(file_path)
                    if text:
                        all_text.append(f"=== {file_name} ===\n{text}")
                        self._log(f"Extracted text ({len(text)} chars)")

                    # Extract embedded images (FAST - these are always available)
                    # Limit to prevent timeout - max 15 embedded images (for 10-25 slide PPTs)
                    img_paths = self._extract_pptx_images(file_path, tmp)
                    embedded_count = len(img_paths)
                    max_embedded = min(15, embedded_count)  # Limit to 15 for speed
                    if embedded_count > max_embedded:
                        self._log(f"Limiting embedded images to {max_embedded} (out of {embedded_count}) to prevent timeout")
                        img_paths = img_paths[:max_embedded]
                    
                    encoded_embedded = 0
                    for p in img_paths:
                        url = self._to_data_url(p, max_size_mb=2.0)  # Allow up to 2MB per image
                        if url:
                            all_images.append({"type": "image_url", "image_url": {"url": url}})
                            encoded_embedded += 1
                    if encoded_embedded > 0:
                        self._log(f"Extracted {encoded_embedded}/{embedded_count} embedded images")

                    # Convert to PDF for page rendering (PRIORITY - do this first for large PPTs)
                    # PDF conversion gives us full slide images which are more important than EMF
                    elapsed = time.time() - start_time
                    time_remaining = self.valves.max_processing_time - elapsed
                    
                    # Calculate time needed: PDF conversion + rendering
                    # For 22 slides: ~30s base + 22*3 = 96s max, but capped at 120s
                    pdf_timeout_needed = min(
                        self.valves.libreoffice_base_timeout + (slide_count * self.valves.libreoffice_per_slide_timeout),
                        self.valves.max_timeout
                    )
                    # More realistic render time estimate: ~2-3s per page for 600 DPI
                    render_time_needed = max(5, slide_count * 2)  # At least 5s, or 2s per slide
                    total_pdf_time_needed = pdf_timeout_needed + render_time_needed + 10  # Add 10s buffer for safety
                    
                    self._log(f"Time check: {int(time_remaining)}s remaining, need ~{int(total_pdf_time_needed)}s for PDF (conversion: {int(pdf_timeout_needed)}s, render: {int(render_time_needed)}s)")
                    
                    pdf_pages_rendered = 0  # Track if PDF conversion succeeded
                    
                    # Try PDF conversion if we have at least 80% of needed time (more lenient)
                    if time_remaining > (total_pdf_time_needed * 0.8):
                        self._log(f"Attempting PDF conversion ({slide_count} slides, {int(time_remaining)}s available)")
                        pdf_path = self._convert_pptx_to_pdf(file_path, tmp)
                        if pdf_path:
                            elapsed = time.time() - start_time
                            time_remaining = self.valves.max_processing_time - elapsed
                            self._log(f"PDF conversion completed in {elapsed:.1f}s, {int(time_remaining)}s remaining for rendering")
                            if time_remaining > 5:  # Need at least 5s for rendering
                                self._log(f"Rendering PDF pages ({int(time_remaining)}s remaining)")
                                page_paths = self._convert_pdf_to_images(pdf_path, tmp)
                                pdf_pages_rendered = len(page_paths)
                                
                                # Limit images to prevent timeout - process max 25 pages (handles typical 10-25 slide PPTs)
                                max_pages_to_encode = min(25, pdf_pages_rendered)
                                if pdf_pages_rendered > max_pages_to_encode:
                                    self._log(f"Limiting to {max_pages_to_encode} pages (out of {pdf_pages_rendered}) to prevent timeout")
                                    page_paths = page_paths[:max_pages_to_encode]
                                
                                # Encode images with timeout check
                                encoded_count = 0
                                for i, p in enumerate(page_paths):
                                    # Check time remaining every 5 images (less frequent checks for speed)
                                    if i > 0 and i % 5 == 0:
                                        elapsed = time.time() - start_time
                                        if elapsed > self.valves.max_processing_time - 15:  # Stop 15s before timeout
                                            self._log(f"Stopping image encoding - {int(self.valves.max_processing_time - elapsed)}s before timeout")
                                            break
                                    
                                    url = self._to_data_url(p, max_size_mb=2.0)  # Allow up to 2MB per image
                                    if url:
                                        all_images.append({"type": "image_url", "image_url": {"url": url}})
                                        encoded_count += 1
                                
                                elapsed_after_render = time.time() - start_time
                                self._log(f"Rendered {encoded_count}/{pdf_pages_rendered} PDF pages (total time: {elapsed_after_render:.1f}s)")
                            else:
                                self._log(f"Skipped PDF rendering - only {int(time_remaining)}s remaining (need 5s+)")
                        else:
                            elapsed = time.time() - start_time
                            self._log(f"PDF conversion failed/timed out after {elapsed:.1f}s - continuing with embedded images")
                    else:
                        self._log(f"Skipped PDF conversion - not enough time ({int(time_remaining)}s remaining, need ~{int(total_pdf_time_needed)}s)")
                        self._log(f"Returning with {embedded_count} embedded images and text")
                    
                    # Extract EMF/WMF (LOWEST PRIORITY - do AFTER PDF conversion)
                    # EMF conversion is slow but useful - do it if we have time after PDF
                    elapsed = time.time() - start_time
                    time_remaining = self.valves.max_processing_time - elapsed
                    
                    # Only do EMF if:
                    # 1. PDF conversion already succeeded (so we have slide images)
                    # 2. We have at least 20s remaining (EMF can be slow)
                    if pdf_pages_rendered > 0 and time_remaining > 20:
                        self._log(f"Attempting EMF/WMF conversion ({int(time_remaining)}s remaining, {pdf_pages_rendered} PDF pages done)")
                        emf_paths = self._extract_emf_wmf(file_path, tmp)
                        emf_count = len(emf_paths)
                        for p in emf_paths:
                            url = self._to_data_url(p)
                            if url:
                                all_images.append({"type": "image_url", "image_url": {"url": url}})
                        if emf_count > 0:
                            self._log(f"Converted {emf_count} EMF/WMF images")
                    else:
                        if pdf_pages_rendered == 0:
                            self._log(f"Skipping EMF conversion - PDF conversion should complete first")
                        else:
                            self._log(f"Skipping EMF conversion - not enough time ({int(time_remaining)}s remaining, need 20s+)")
                    
                    elapsed = time.time() - start_time
                    self._log(f"Total processing time: {elapsed:.1f}s, {len(all_images)} total images")

                # PDF processing
                elif is_pdf:
                    page_paths = self._convert_pdf_to_images(file_path, tmp)
                    for p in page_paths:
                        url = self._to_data_url(p)
                        if url:
                            all_images.append({"type": "image_url", "image_url": {"url": url}})

            # Mark as processed
            self._processed_files.add(file_hash)

        # Build response
        if not all_text and not all_images:
            self._log("No content extracted")
            return body

        last_msg = messages[-1]
        original = last_msg.get("content", "")
        if isinstance(original, list):
            original = " ".join(
                item.get("text", "") for item in original 
                if isinstance(item, dict) and item.get("type") == "text"
            )

        # Build content blocks
        text_content = original.strip()
        if all_text:
            text_content += "\n\n--- Extracted Document Content ---\n" + "\n\n".join(all_text)
        if all_images:
            text_content += f"\n\n[{len(all_images)} document images attached for visual analysis]"

        content_blocks = [{"type": "text", "text": text_content}]
        for img in all_images:
            if img.get("image_url", {}).get("url"):
                content_blocks.append(img)

        messages[-1]["content"] = content_blocks
        body["messages"] = messages

        self._log(f"Done: {len(all_text)} text sections, {len(all_images)} images")
        self._log("=" * 60)
        
        return body

    def stream(self, event: dict, __user__: Optional[dict] = None) -> dict:
        return event

    def outlet(self, body: dict, __user__: Optional[dict] = None) -> dict:
        return body
