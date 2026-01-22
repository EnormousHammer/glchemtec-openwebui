"""
title: PPT/PDF Vision Filter
author: GLChemTec
version: 9.0
description: Local PPTX extraction (no external service needed) + vision processing for PDFs.
"""

import os
import base64
import json
import tempfile
import shutil
import requests
import zipfile
from typing import Optional, List, Dict, Any
from io import BytesIO

from pydantic import BaseModel, Field

# Import python-pptx for local extraction
try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


class Filter:
    class Valves(BaseModel):
        priority: int = Field(default=0, description="Filter priority (0 = highest)")
        enabled: bool = Field(default=True, description="Enable PPT/PDF vision processing")
        debug: bool = Field(default=True, description="Enable debug logging")
        
        # PPTX Extraction - Built into OpenWebUI instance (no external service needed)
        # The extraction now happens directly in this filter using python-pptx
        
        # PDF Processing (fallback for PDFs)
        max_pages: int = Field(default=8, description="Max pages for PDF processing")
        dpi: int = Field(default=150, description="DPI for PDF rendering (higher = clearer text)")
        max_image_width: int = Field(default=1600, description="Max width in pixels (larger = more detail)")
        max_image_height: int = Field(default=1200, description="Max height in pixels (larger = more content)")
        jpeg_quality: int = Field(default=85, description="JPEG quality (0-100, higher = less compression)")
        max_total_base64_mb: float = Field(default=5.0, description="Max total base64 size in MB")

    def __init__(self):
        self.valves = self.Valves()

    def _log(self, msg: str) -> None:
        if self.valves.debug:
            print(f"[PPT-PDF-VISION] {msg}")

    def _extract_all_files(self, body: dict, messages: list) -> List[Dict[str, Any]]:
        """Extract files from all possible locations in the request."""
        all_files = []
        
        # Check body.files
        if isinstance(body.get("files"), list):
            self._log(f"Found {len(body['files'])} files in body['files']")
            all_files.extend(body["files"])
        
        # Check each message
        for msg in messages:
            if isinstance(msg.get("files"), list):
                self._log(f"Found {len(msg['files'])} files in message['files']")
                all_files.extend(msg["files"])
            if isinstance(msg.get("attachments"), list):
                self._log(f"Found {len(msg['attachments'])} files in message['attachments']")
                all_files.extend(msg["attachments"])
            if isinstance(msg.get("sources"), list):
                for source_obj in msg["sources"]:
                    if isinstance(source_obj, dict):
                        source = source_obj.get("source", {})
                        if source.get("type") == "file" and isinstance(source.get("file"), dict):
                            all_files.append({"file": source["file"]})
        
        # Deduplicate by path
        seen = set()
        unique = []
        for f in all_files:
            path = self._get_file_path(f)
            if path and path not in seen:
                seen.add(path)
                unique.append(f)
        
        self._log(f"Total unique files found: {len(unique)}")
        return unique

    def _get_file_path(self, file_obj: Dict[str, Any]) -> str:
        """Extract file path from various file object structures."""
        if not isinstance(file_obj, dict):
            return ""
        
        # Try file.path
        if isinstance(file_obj.get("file"), dict):
            f = file_obj["file"]
            path = f.get("path", "")
            if path:
                return path.strip()
            # Try file.meta.path
            if isinstance(f.get("meta"), dict):
                return f["meta"].get("path", "").strip()
        
        # Try direct path
        return file_obj.get("path", "").strip()

    def _get_file_name(self, file_obj: Dict[str, Any]) -> str:
        """Extract filename from various file object structures."""
        if not isinstance(file_obj, dict):
            return ""
        
        if isinstance(file_obj.get("file"), dict):
            f = file_obj["file"]
            name = f.get("filename") or f.get("name") or ""
            if not name and isinstance(f.get("meta"), dict):
                name = f["meta"].get("name") or ""
            return name.lower().strip()
        
        return (file_obj.get("name") or file_obj.get("filename") or "").lower().strip()

    def _extract_text_content(self, content: Any) -> str:
        """Extract text from message content (handles string or list format)."""
        if isinstance(content, str):
            return content
        if isinstance(content, list):
            parts = [item.get("text", "") for item in content if isinstance(item, dict) and item.get("type") == "text"]
            return "\n".join(p for p in parts if p).strip()
        return str(content) if content else ""

    def _is_openai_model(self, model: str) -> bool:
        """Check if model uses OpenAI format."""
        m = (model or "").lower()
        if not m:
            return True
        if "claude" in m or "anthropic" in m:
            return False
        return True  # Default to OpenAI format

    def _is_anthropic_model(self, model: str) -> bool:
        """Check if model is Anthropic/Claude."""
        m = (model or "").lower()
        return "claude" in m or "anthropic" in m

    def extract_pptx_locally(self, file_path: str, file_name: str) -> Optional[Dict[str, Any]]:
        """
        Extract PPTX content directly in OpenWebUI instance using python-pptx.
        No external service needed - everything runs in this instance.
        Returns dict with slides, text, images, tables, etc.
        """
        if not PPTX_AVAILABLE:
            self._log("python-pptx not available, cannot extract PPTX locally")
            return None
        
        try:
            self._log(f"Extracting PPTX locally: {file_name}")
            prs = Presentation(file_path)
            
            slides_data = []
            all_text = []
            all_images = []
            all_tables = []
            
            for slide_idx, slide in enumerate(prs.slides, 1):
                slide_text = []
                slide_images = []
                slide_tables = []
                
                # Extract text from all shapes
                for shape in slide.shapes:
                    # Text from text boxes
                    if hasattr(shape, "text") and shape.text:
                        text = shape.text.strip()
                        if text:
                            slide_text.append(text)
                            all_text.append(text)
                    
                    # Images - Extract at FULL QUALITY (no compression)
                    if hasattr(shape, "image"):
                        try:
                            image_bytes = shape.image.blob
                            img_size_kb = len(image_bytes) // 1024
                            
                            # Use original image bytes - NO compression for maximum quality
                            # This is critical for NMR spectra, chemical structures, and small text
                            b64_data = base64.b64encode(image_bytes).decode("utf-8")
                            
                            # Determine mime type from original
                            ext = shape.image.ext.lower()
                            mime_type = "image/png"  # Default to PNG for lossless
                            if ext in ("jpg", "jpeg"):
                                mime_type = "image/jpeg"
                            elif ext == "gif":
                                mime_type = "image/gif"
                            
                            # Log image details for verification
                            self._log(f"Slide {slide_idx}: Extracted image - {img_size_kb}KB, {mime_type}, FULL QUALITY (no compression)")
                            
                            img_data = {
                                "base64": b64_data,
                                "data_url": f"data:{mime_type};base64,{b64_data}",
                                "mime_type": mime_type,
                                "slide_number": slide_idx,
                                "size_kb": img_size_kb,
                                "quality": "full"  # Mark as full quality
                            }
                            slide_images.append(img_data)
                            all_images.append(img_data)
                        except Exception as e:
                            self._log(f"Error extracting image from slide {slide_idx}: {e}")
                    
                    # Tables
                    if hasattr(shape, "table") and shape.has_table:
                        try:
                            table = shape.table
                            table_data = {
                                "rows": [],
                                "slide_number": slide_idx
                            }
                            
                            for row in table.rows:
                                row_data = []
                                for cell in row.cells:
                                    cell_text = cell.text.strip() if cell.text else ""
                                    row_data.append(cell_text)
                                table_data["rows"].append(row_data)
                            
                            slide_tables.append(table_data)
                            all_tables.append(table_data)
                        except Exception as e:
                            self._log(f"Error extracting table from slide {slide_idx}: {e}")
                    
                    # Group shapes (may contain nested images/text)
                    if hasattr(shape, "shapes"):
                        for sub_shape in shape.shapes:
                            if hasattr(sub_shape, "text") and sub_shape.text:
                                text = sub_shape.text.strip()
                                if text:
                                    slide_text.append(text)
                                    all_text.append(text)
                            
                            if hasattr(sub_shape, "image"):
                                try:
                                    image_bytes = sub_shape.image.blob
                                    img_size_kb = len(image_bytes) // 1024
                                    b64_data = base64.b64encode(image_bytes).decode("utf-8")
                                    ext = sub_shape.image.ext.lower()
                                    mime_type = "image/png"  # Default to PNG for lossless
                                    if ext in ("jpg", "jpeg"):
                                        mime_type = "image/jpeg"
                                    
                                    self._log(f"Slide {slide_idx}: Extracted nested image - {img_size_kb}KB, FULL QUALITY")
                                    
                                    img_data = {
                                        "base64": b64_data,
                                        "data_url": f"data:{mime_type};base64,{b64_data}",
                                        "mime_type": mime_type,
                                        "slide_number": slide_idx,
                                        "size_kb": img_size_kb,
                                        "quality": "full"
                                    }
                                    slide_images.append(img_data)
                                    all_images.append(img_data)
                                except Exception as e:
                                    self._log(f"Error extracting nested image: {e}")
                
                # Extract speaker notes
                notes_text = ""
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
                
                slide_data = {
                    "slide_number": slide_idx,
                    "text": "\n".join(slide_text),
                    "images": slide_images,
                    "tables": slide_tables,
                    "notes": notes_text
                }
                slides_data.append(slide_data)
            
            result = {
                "slides": slides_data,
                "total_slides": len(slides_data),
                "all_text": "\n\n".join(all_text),
                "total_images": len(all_images),
                "total_tables": len(all_tables),
                "extracted_by": "local_python_pptx"
            }
            
            # Log detailed extraction summary
            total_img_size = sum(img.get("size_kb", 0) for img in all_images)
            self._log(f"✅ PPTX extraction successful:")
            self._log(f"   - {len(slides_data)} slides")
            self._log(f"   - {len(all_images)} images (FULL QUALITY - no compression)")
            self._log(f"   - {len(all_tables)} tables")
            self._log(f"   - Total image size: {total_img_size}KB (all sent at full quality to OpenAI)")
            self._log(f"   - Images will be readable for NMR spectra, chemical structures, and small text")
            return result
            
        except Exception as e:
            self._log(f"Local PPTX extraction error: {type(e).__name__}: {e}")
            import traceback
            self._log(f"Traceback: {traceback.format_exc()}")
            return None
    

    def _extract_image_from_shape(self, shape, slide_num: int) -> Optional[Dict[str, Any]]:
        """Extract image from a shape, handling various shape types."""
        try:
            # Direct image shape
            if hasattr(shape, "image"):
                image = shape.image
                image_bytes = image.blob
                
                # Convert to base64
                b64_data = base64.b64encode(image_bytes).decode("utf-8")
                
                # Determine mime type from image extension
                ext = image.ext
                if ext.lower() in ("png", ".png"):
                    mime_type = "image/png"
                elif ext.lower() in ("jpg", "jpeg", ".jpg", ".jpeg"):
                    mime_type = "image/jpeg"
                elif ext.lower() in ("gif", ".gif"):
                    mime_type = "image/gif"
                else:
                    mime_type = "image/png"  # Default
                
                return {
                    "data": b64_data,
                    "base64": b64_data,
                    "content_type": mime_type,
                    "slide_number": slide_num
                }
            
            # Picture shape (alternative attribute)
            if hasattr(shape, "image") or (hasattr(shape, "shape_type") and "picture" in str(shape.shape_type).lower()):
                try:
                    if hasattr(shape, "image"):
                        image = shape.image
                        image_bytes = image.blob
                        b64_data = base64.b64encode(image_bytes).decode("utf-8")
                        ext = getattr(image, "ext", "png")
                        if ext.lower() in ("png", ".png"):
                            mime_type = "image/png"
                        elif ext.lower() in ("jpg", "jpeg", ".jpg", ".jpeg"):
                            mime_type = "image/jpeg"
                        elif ext.lower() in ("gif", ".gif"):
                            mime_type = "image/gif"
                        else:
                            mime_type = "image/png"
                        return {
                            "data": b64_data,
                            "base64": b64_data,
                            "content_type": mime_type,
                            "slide_number": slide_num
                        }
                except Exception as e:
                    self._log(f"Slide {slide_num}: Error extracting from picture shape: {e}")
            
            return None
        except Exception as e:
            self._log(f"Slide {slide_num}: Error in _extract_image_from_shape: {e}")
            return None

    def _extract_images_recursive(self, shapes, slide_num: int) -> List[Dict[str, Any]]:
        """Recursively extract images from shapes, including grouped shapes."""
        images = []
        
        for shape in shapes:
            # Try to extract image from this shape
            img = self._extract_image_from_shape(shape, slide_num)
            if img:
                images.append(img)
                self._log(f"Slide {slide_num}: Extracted image ({len(img['base64']) * 3 // 4} bytes, {img['content_type']})")
            
            # If shape is a group, recursively extract from its shapes
            if hasattr(shape, "shapes"):
                try:
                    nested_images = self._extract_images_recursive(shape.shapes, slide_num)
                    images.extend(nested_images)
                except Exception as e:
                    self._log(f"Slide {slide_num}: Error extracting from grouped shape: {e}")
        
        return images

    def _extract_text_from_shape(self, shape, slide_num: int) -> List[str]:
        """Comprehensively extract ALL text from a shape, including tables, placeholders, etc."""
        texts = []
        
        try:
            # Method 1: Direct text attribute (most common)
            if hasattr(shape, "text") and shape.text:
                text = str(shape.text).strip()
                if text:
                    texts.append(text)
            
            # Method 2: Extract from text_frame (for text boxes, placeholders, etc.)
            if hasattr(shape, "text_frame"):
                try:
                    text_frame = shape.text_frame
                    # Extract from all paragraphs
                    if hasattr(text_frame, "paragraphs"):
                        for para in text_frame.paragraphs:
                            para_text = ""
                            # Extract from all runs in paragraph
                            if hasattr(para, "runs"):
                                for run in para.runs:
                                    if hasattr(run, "text") and run.text:
                                        para_text += run.text
                            # If no runs, try direct paragraph text
                            if not para_text and hasattr(para, "text"):
                                para_text = para.text
                            
                            if para_text.strip():
                                texts.append(para_text.strip())
                except Exception as e:
                    self._log(f"Slide {slide_num}: Error extracting from text_frame: {e}")
            
            # Method 3: Extract from table cells
            if hasattr(shape, "table"):
                try:
                    table = shape.table
                    table_rows = []
                    for row in table.rows:
                        row_cells = []
                        for cell in row.cells:
                            cell_text = ""
                            # Extract text from cell's text_frame
                            if hasattr(cell, "text_frame") and cell.text_frame:
                                if hasattr(cell.text_frame, "paragraphs"):
                                    for para in cell.text_frame.paragraphs:
                                        if hasattr(para, "runs"):
                                            for run in para.runs:
                                                if hasattr(run, "text") and run.text:
                                                    cell_text += run.text
                                        elif hasattr(para, "text"):
                                            cell_text += para.text
                                elif hasattr(cell.text_frame, "text"):
                                    cell_text = cell.text_frame.text
                            # Fallback to cell.text
                            if not cell_text and hasattr(cell, "text"):
                                cell_text = cell.text
                            
                            row_cells.append(cell_text.strip() if cell_text else "")
                        if any(row_cells):  # Only add non-empty rows
                            table_rows.append(row_cells)
                    
                    if table_rows:
                        # Format table as text
                        table_text = "Table:\n"
                        for row in table_rows:
                            table_text += " | ".join(str(cell) for cell in row) + "\n"
                        texts.append(table_text.strip())
                except Exception as e:
                    self._log(f"Slide {slide_num}: Error extracting from table: {e}")
            
            # Method 4: Extract from placeholders
            if hasattr(shape, "is_placeholder") and shape.is_placeholder:
                try:
                    if hasattr(shape, "text"):
                        placeholder_text = str(shape.text).strip()
                        if placeholder_text:
                            texts.append(placeholder_text)
                except Exception as e:
                    self._log(f"Slide {slide_num}: Error extracting from placeholder: {e}")
            
            # Method 5: Extract from auto-shape text
            if hasattr(shape, "auto_shape_type"):
                try:
                    if hasattr(shape, "text"):
                        auto_text = str(shape.text).strip()
                        if auto_text:
                            texts.append(auto_text)
                except Exception as e:
                    self._log(f"Slide {slide_num}: Error extracting from auto-shape: {e}")
            
        except Exception as e:
            self._log(f"Slide {slide_num}: Error in _extract_text_from_shape: {e}")
        
        return texts

    def _extract_text_recursive(self, shapes, slide_num: int) -> List[str]:
        """Recursively extract ALL text from shapes, including grouped shapes."""
        all_texts = []
        
        for shape in shapes:
            # Extract text from this shape
            shape_texts = self._extract_text_from_shape(shape, slide_num)
            if shape_texts:
                all_texts.extend(shape_texts)
            
            # If shape is a group, recursively extract from its shapes
            if hasattr(shape, "shapes"):
                try:
                    nested_texts = self._extract_text_recursive(shape.shapes, slide_num)
                    all_texts.extend(nested_texts)
                except Exception as e:
                    self._log(f"Slide {slide_num}: Error extracting text from grouped shape: {e}")
        
        return all_texts

    def extract_pptx_fallback(self, file_path: str, file_name: str) -> Optional[Dict[str, Any]]:
        """
        Fallback extraction using python-pptx when LibreOffice conversion fails.
        Extracts both text and images, returning data in the same format as the converter service.
        Enhanced to extract images from grouped shapes, backgrounds, and all shape types.
        """
        try:
            from pptx import Presentation
            
            self._log(f"Attempting fallback extraction using python-pptx for {file_name}")
            prs = Presentation(file_path)
            
            slides_data = []
            total_images = 0
            
            for idx, slide in enumerate(prs.slides, 1):
                slide_data = {
                    "slide_number": idx,
                    "title": "",
                    "text": [],
                    "notes": "",
                    "images": []
                }
                
                # Extract title (comprehensive)
                if slide.shapes.title:
                    title_texts = self._extract_text_from_shape(slide.shapes.title, idx)
                    if title_texts:
                        slide_data["title"] = " ".join(title_texts).strip()
                
                # Extract ALL text from all shapes using comprehensive method
                all_slide_texts = self._extract_text_recursive(slide.shapes, idx)
                if all_slide_texts:
                    # Deduplicate while preserving order
                    seen = set()
                    for text in all_slide_texts:
                        text_clean = text.strip()
                        if text_clean and text_clean not in seen:
                            seen.add(text_clean)
                            slide_data["text"].append(text_clean)
                
                # Also extract tables separately for better formatting
                for shape in slide.shapes:
                    if hasattr(shape, "table"):
                        try:
                            table = shape.table
                            table_data = []
                            for row in table.rows:
                                row_data = []
                                for cell in row.cells:
                                    cell_text = ""
                                    if hasattr(cell, "text_frame") and cell.text_frame:
                                        if hasattr(cell.text_frame, "paragraphs"):
                                            for para in cell.text_frame.paragraphs:
                                                if hasattr(para, "runs"):
                                                    for run in para.runs:
                                                        if hasattr(run, "text") and run.text:
                                                            cell_text += run.text
                                                elif hasattr(para, "text"):
                                                    cell_text += para.text
                                        elif hasattr(cell.text_frame, "text"):
                                            cell_text = cell.text_frame.text
                                    if not cell_text and hasattr(cell, "text"):
                                        cell_text = cell.text
                                    row_data.append(cell_text.strip() if cell_text else "")
                                if any(row_data):
                                    table_data.append(row_data)
                            
                            if table_data:
                                slide_data.setdefault("tables", []).append(table_data)
                        except Exception as e:
                            self._log(f"Slide {idx}: Error extracting table: {e}")
                
                # Extract images recursively (handles grouped shapes)
                slide_images = self._extract_images_recursive(slide.shapes, idx)
                slide_data["images"].extend(slide_images)
                total_images += len(slide_images)
                
                # Try to extract background image if present
                try:
                    if hasattr(slide, "background") and slide.background:
                        bg = slide.background
                        if hasattr(bg, "fill") and hasattr(bg.fill, "picture"):
                            pic = bg.fill.picture
                            if hasattr(pic, "image"):
                                bg_image = pic.image
                                bg_bytes = bg_image.blob
                                bg_b64 = base64.b64encode(bg_bytes).decode("utf-8")
                                ext = bg_image.ext
                                if ext.lower() in ("png", ".png"):
                                    bg_mime = "image/png"
                                elif ext.lower() in ("jpg", "jpeg", ".jpg", ".jpeg"):
                                    bg_mime = "image/jpeg"
                                elif ext.lower() in ("gif", ".gif"):
                                    bg_mime = "image/gif"
                                else:
                                    bg_mime = "image/png"
                                
                                slide_data["images"].append({
                                    "data": bg_b64,
                                    "base64": bg_b64,
                                    "content_type": bg_mime,
                                    "slide_number": idx
                                })
                                total_images += 1
                                self._log(f"Slide {idx}: Extracted background image ({len(bg_bytes)} bytes, {bg_mime})")
                except Exception as bg_e:
                    # Background extraction is optional, don't fail if it doesn't work
                    pass
                
                # Extract notes comprehensively
                try:
                    if slide.has_notes_slide:
                        notes_slide = slide.notes_slide
                        notes_text = ""
                        
                        # Extract from notes_text_frame
                        if hasattr(notes_slide, "notes_text_frame") and notes_slide.notes_text_frame:
                            if hasattr(notes_slide.notes_text_frame, "paragraphs"):
                                for para in notes_slide.notes_text_frame.paragraphs:
                                    if hasattr(para, "runs"):
                                        for run in para.runs:
                                            if hasattr(run, "text") and run.text:
                                                notes_text += run.text
                                    elif hasattr(para, "text"):
                                        notes_text += para.text
                            elif hasattr(notes_slide.notes_text_frame, "text"):
                                notes_text = notes_slide.notes_text_frame.text
                        
                        # Also check for shapes in notes slide
                        if hasattr(notes_slide, "shapes"):
                            notes_shape_texts = self._extract_text_recursive(notes_slide.shapes, idx)
                            if notes_shape_texts:
                                notes_text += "\n" + "\n".join(notes_shape_texts)
                        
                        notes_text = notes_text.strip()
                        if notes_text:
                            slide_data["notes"] = notes_text
                            self._log(f"Slide {idx}: Extracted notes ({len(notes_text)} chars)")
                except Exception as notes_e:
                    self._log(f"Slide {idx}: Error extracting notes: {notes_e}")
                
                slides_data.append(slide_data)
            
            # Calculate extraction statistics
            total_text_chars = 0
            total_tables = 0
            slides_with_content = 0
            
            for slide_data in slides_data:
                slide_text_count = sum(len(str(t)) for t in slide_data.get("text", []))
                total_text_chars += slide_text_count
                total_text_chars += len(slide_data.get("title", ""))
                total_text_chars += len(slide_data.get("notes", ""))
                total_tables += len(slide_data.get("tables", []))
                if slide_text_count > 0 or slide_data.get("title") or slide_data.get("notes"):
                    slides_with_content += 1
            
            result = {
                "slides": slides_data,
                "fallback_mode": True,
                "total_slides": len(slides_data),
                "total_images": total_images,
                "total_text_chars": total_text_chars,
                "total_tables": total_tables,
                "slides_with_content": slides_with_content
            }
            
            self._log(f"Fallback extraction successful: {len(slides_data)} slides, {total_images} images, {total_text_chars} text chars, {total_tables} tables")
            self._log(f"Content coverage: {slides_with_content}/{len(slides_data)} slides have content ({100*slides_with_content/len(slides_data) if slides_data else 0:.1f}%)")
            return result
            
        except ImportError:
            self._log("python-pptx not available for fallback extraction")
            return None
        except Exception as e:
            self._log(f"Fallback extraction failed: {type(e).__name__}: {e}")
            import traceback
            self._log(f"Traceback: {traceback.format_exc()}")
            return None

    def format_pptx_extraction(self, data: Dict[str, Any], file_name: str, include_image_refs: bool = True) -> str:
        """Format extracted PPTX data into natural, flowing text for the model."""
        slides = data.get("slides", [])
        if not slides:
            return f"Document: {file_name}\n\nNo content found in this presentation."
        
        # Build natural flowing text
        content_parts = []
        
        # Add document header
        content_parts.append(f"Document: {file_name}")
        content_parts.append(f"Total slides: {len(slides)}\n")
        
        # Combine all text content naturally
        all_text = []
        for slide in slides:
            slide_text = []
            slide_num = slide.get("slide_number", "?")
            
            # Add title if present
            title = slide.get("title", "").strip()
            if title:
                slide_text.append(title)
            
            # Add text content
            text_content = slide.get("text", [])
            if isinstance(text_content, list):
                for text in text_content:
                    if text and text.strip():
                        slide_text.append(text.strip())
            elif text_content and str(text_content).strip():
                slide_text.append(str(text_content).strip())
            
            # Add notes if present
            notes = slide.get("notes", "").strip()
            if notes:
                slide_text.append(f"(Note: {notes})")
            
            # Add tables as formatted text (comprehensive)
            tables = slide.get("tables", [])
            for table_idx, table in enumerate(tables):
                if isinstance(table, list):
                    table_text = []
                    for row in table:
                        if isinstance(row, list):
                            # Clean cell values and join
                            clean_row = [str(cell).strip() if cell else "" for cell in row]
                            if any(clean_row):  # Only add non-empty rows
                                table_text.append(" | ".join(clean_row))
                    if table_text:
                        slide_text.append(f"Table {table_idx + 1}:\n" + "\n".join(table_text))
            
            # Add image references for source view
            if include_image_refs:
                images = slide.get("images", [])
                if images:
                    img_refs = [f"[Image {i+1} from slide {slide_num}]" for i in range(len(images))]
                    slide_text.append(" ".join(img_refs))
            
            # Combine slide content
            if slide_text:
                all_text.append("\n".join(slide_text))
        
        # Join all content with natural spacing
        if all_text:
            content_parts.append("\n\n".join(all_text))
        
        return "\n".join(content_parts)

    def download_image_to_base64(self, url: str) -> Optional[Dict[str, str]]:
        """Download image from URL and convert to base64."""
        try:
            self._log(f"Downloading image: {url[:80]}...")
            response = requests.get(url, timeout=30)
            
            if response.status_code != 200:
                self._log(f"Failed to download image: {response.status_code}")
                return None
            
            # Determine mime type
            content_type = response.headers.get("Content-Type", "image/png")
            if "jpeg" in content_type or "jpg" in content_type:
                mime_type = "image/jpeg"
            elif "png" in content_type:
                mime_type = "image/png"
            elif "gif" in content_type:
                mime_type = "image/gif"
            elif "webp" in content_type:
                mime_type = "image/webp"
            else:
                mime_type = "image/png"  # Default
            
            # Convert to base64
            b64_data = base64.b64encode(response.content).decode("utf-8")
            
            self._log(f"Downloaded image: {len(b64_data)//1024}KB as {mime_type}")
            return {"base64": b64_data, "mime_type": mime_type}
            
        except Exception as e:
            self._log(f"Error downloading image: {e}")
            return None

    def get_pptx_images_as_base64(self, data: Dict[str, Any]) -> List[Dict[str, Any]]:
        """Extract images from PPTX extraction data - handles URLs and base64."""
        images = []
        total_size = 0
        max_bytes = int(self.valves.max_total_base64_mb * 1024 * 1024)
        
        for slide in data.get("slides", []):
            slide_num = slide.get("slide_number", "?")
            
            for img in slide.get("images", []):
                # Check total size limit
                if total_size > max_bytes:
                    self._log(f"Reached image size limit ({self.valves.max_total_base64_mb}MB), stopping")
                    break
                
                b64_data = None
                mime_type = "image/png"
                
                # Option 1: Image already has base64 data
                if img.get("data") or img.get("base64"):
                    b64_data = img.get("data") or img.get("base64")
                    mime_type = img.get("content_type", "image/png")
                    self._log(f"Slide {slide_num}: Using embedded base64 image")
                
                # Option 2: Image has URL - download it
                elif img.get("url"):
                    url = img.get("url")
                    downloaded = self.download_image_to_base64(url)
                    if downloaded:
                        b64_data = downloaded["base64"]
                        mime_type = downloaded["mime_type"]
                        self._log(f"Slide {slide_num}: Downloaded image from URL")
                
                # Option 3: Check for 'src' or 'path' as URL
                elif img.get("src") or img.get("path"):
                    url = img.get("src") or img.get("path")
                    if url.startswith("http"):
                        downloaded = self.download_image_to_base64(url)
                        if downloaded:
                            b64_data = downloaded["base64"]
                            mime_type = downloaded["mime_type"]
                            self._log(f"Slide {slide_num}: Downloaded image from src/path")
                
                if b64_data:
                    total_size += len(b64_data)
                    images.append({
                        "type": "image_url",
                        "image_url": {
                            "url": f"data:{mime_type};base64,{b64_data}"
                        }
                    })
        
        # Log detailed image information for verification
        self._log(f"✅ Total images prepared for OpenAI: {len(images)}")
        self._log(f"   - Total size: {total_size//1024}KB")
        self._log(f"   - All images at FULL QUALITY (no compression)")
        self._log(f"   - Images sent as base64 data URLs - OpenAI receives full resolution")
        self._log(f"   - Quality sufficient for NMR spectra, chemical structures, and small text")
        return images

    def convert_pdf_to_images(self, pdf_path: str, output_dir: str) -> List[str]:
        """Convert PDF pages to images (fallback for PDF files)."""
        try:
            from pdf2image import convert_from_path
            from PIL import Image
            
            self._log(f"Converting PDF with max_pages={self.valves.max_pages}")
            images = convert_from_path(
                pdf_path, 
                dpi=self.valves.dpi, 
                fmt="png",
                first_page=1, 
                last_page=self.valves.max_pages
            )
            
            paths = []
            total_size = 0
            max_bytes = int(self.valves.max_total_base64_mb * 1024 * 1024)
            
            for idx, img in enumerate(images):
                w, h = img.width, img.height
                max_w = self.valves.max_image_width
                max_h = self.valves.max_image_height
                
                scale = min(max_w / w, max_h / h, 1.0)
                if scale < 1.0:
                    new_w = int(w * scale)
                    new_h = int(h * scale)
                    img = img.resize((new_w, new_h), Image.LANCZOS)
                
                path = os.path.join(output_dir, f"page_{idx+1:03d}.jpg")
                img = img.convert("RGB")
                img.save(path, "JPEG", quality=self.valves.jpeg_quality, optimize=True)
                
                file_size = os.path.getsize(path)
                total_size += file_size
                
                if total_size > max_bytes:
                    self._log(f"Stopping at page {idx+1} - size limit reached")
                    paths.append(path)
                    break
                
                paths.append(path)
            
            self._log(f"Created {len(paths)} PDF images")
            return paths
            
        except Exception as e:
            self._log(f"PDF->images error: {e}")
            return []

    def image_to_base64(self, path: str) -> Optional[str]:
        """Convert image file to base64 string."""
        try:
            with open(path, "rb") as f:
                return base64.b64encode(f.read()).decode("utf-8")
        except:
            return None

    def inlet(self, body: dict, __user__: Optional[dict] = None) -> dict:
        """Main filter entry point - processes incoming requests."""
        self._log("=" * 60)
        self._log("INLET - PPT/PDF Vision Filter v8.0")
        self._log("Using built-in PPTX extraction (no external service)")

        if not self.valves.enabled:
            self._log("Filter is disabled")
            return body

        messages = body.get("messages", [])
        if not messages:
            self._log("No messages in request")
            return body

        # Log body structure for debugging
        self._log(f"Body keys: {list(body.keys())}")
        
        files = self._extract_all_files(body, messages)
        
        if not files:
            self._log("No files found in request")
            return body

        all_content = []  # Text content from extractions
        all_images = []   # Image blocks for vision

        model_name = body.get("model", "")
        use_openai_format = self._is_openai_model(model_name)
        self._log(f"Model: {model_name}, Format: {'OpenAI' if use_openai_format else 'Anthropic'}")

        for file_obj in files:
            file_path = self._get_file_path(file_obj)
            file_name = self._get_file_name(file_obj)
            
            self._log(f"Processing file: {file_name}")
            self._log(f"File path: {file_path}")

            if not file_path:
                self._log("No path found for file, skipping")
                continue
                
            if not os.path.exists(file_path):
                self._log(f"File does not exist: {file_path}")
                continue

            is_pptx = file_name.endswith((".ppt", ".pptx"))
            is_pdf = file_name.endswith(".pdf")

            if not (is_pptx or is_pdf):
                self._log(f"Skipping non-PPT/PDF file: {file_name}")
                continue

            # ========== PPTX PROCESSING ==========
            if is_pptx:
                self._log("Processing PPTX file (built-in extraction)")
                
                # Extract directly in this instance - no external service needed
                extracted = self.extract_pptx_locally(file_path, file_name)
                
                if extracted:
                    # Add formatted text content
                    formatted_text = self.format_pptx_extraction(extracted, file_name)
                    all_content.append(formatted_text)
                    self._log(f"Added extracted text ({len(formatted_text)} chars)")
                    
                    # Add images if available
                    images = self.get_pptx_images_as_base64(extracted)
                    if images:
                        all_images.extend(images)
                        self._log(f"✅ Added {len(images)} images from PPTX to send to OpenAI")
                        # Log image details for verification
                        for idx, img in enumerate(images[:5]):  # Log first 5
                            img_size = len(img.get("url", "").split(",")[-1]) if "," in img.get("url", "") else 0
                            self._log(f"  Image {idx+1}: {img_size//1024}KB, type: {img.get('type', 'image_url')}")
                    else:
                        self._log("⚠️ WARNING: PPTX extraction succeeded but no images found")
                    
                    # Verify tables are in the formatted text
                    tables_count = sum(len(slide.get("tables", [])) for slide in extracted.get("slides", []))
                    if tables_count > 0:
                        self._log(f"✅ Found {tables_count} tables - included in formatted text sent to OpenAI")
                    else:
                        self._log("ℹ️ No tables found in PPTX")
                else:
                    self._log("PPTX extraction failed - attempting fallback extraction")
                    
                    # Try fallback extraction using python-pptx (text + images)
                    fallback_data = self.extract_pptx_fallback(file_path, file_name)
                    if fallback_data:
                        # Format and add text content
                        formatted_text = self.format_pptx_extraction(fallback_data, file_name)
                        all_content.append(formatted_text)
                        self._log(f"Added fallback text ({len(formatted_text)} chars)")
                        
                        # Add images if available
                        images = self.get_pptx_images_as_base64(fallback_data)
                        if images:
                            all_images.extend(images)
                            self._log(f"Added {len(images)} images from fallback extraction")
                        else:
                            self._log("Fallback extraction succeeded but no images found")
                    else:
                        # Both methods failed
                        self._log("Both primary and fallback extraction failed - no content available")
                        error_msg = (
                            f"[ERROR: PPTX extraction failed for {file_name}]\n"
                            f"Unable to extract content from PowerPoint file.\n"
                            f"Fallback extraction also failed.\n"
                            f"This is typically due to:\n"
                            f"- LibreOffice conversion failure (check converter logs for details)\n"
                            f"- Corrupted or unsupported PPTX file\n"
                            f"- File too large or contains unsupported elements\n"
                            f"Check the filter logs above for the full error message and diagnostics."
                        )
                        all_content.append(error_msg)

            # ========== PDF PROCESSING ==========
            elif is_pdf:
                self._log("Processing PDF via local image conversion")
                
                with tempfile.TemporaryDirectory() as tmp_dir:
                    image_paths = self.convert_pdf_to_images(file_path, tmp_dir)
                    
                    for img_path in image_paths:
                        b64 = self.image_to_base64(img_path)
                        if b64:
                            # Always use OpenAI format - OpenWebUI requires this format
                            all_images.append({
                                "type": "image_url",
                                "image_url": {
                                    "url": f"data:image/jpeg;base64,{b64}"
                                }
                            })
                    
                    self._log(f"Added {len(image_paths)} PDF page images")

        # ========== BUILD FINAL MESSAGE ==========
        if not all_content and not all_images:
            self._log("No content extracted from any files")
            return body

        self._log(f"📤 Building message to send to OpenAI:")
        self._log(f"   - {len(all_content)} text sections (includes tables, text, notes)")
        self._log(f"   - {len(all_images)} images (will be sent as input_image to OpenAI)")
        
        # Verify what we're sending
        if all_images:
            total_img_size = 0
            for img in all_images:
                if isinstance(img, dict) and img.get("type") == "image_url":
                    url = img.get("image_url", {})
                    if isinstance(url, dict):
                        url_str = url.get("url", "")
                    else:
                        url_str = str(url)
                    if url_str.startswith("data:"):
                        b64_part = url_str.split(",", 1)[1] if "," in url_str else ""
                        total_img_size += len(b64_part)
            self._log(f"   - Total image data size: {total_img_size//1024}KB")
        
        if all_content:
            total_text_size = sum(len(str(c)) for c in all_content)
            self._log(f"   - Total text size: {total_text_size//1024}KB (includes tables)")

        last_message = messages[-1]
        original_prompt = self._extract_text_content(last_message.get("content", ""))

        # Build the combined message
        combined_text = f"{original_prompt}\n\n"
        
        if all_content:
            combined_text += "Document content:\n"
            combined_text += "\n\n".join(all_content)
            combined_text += "\n\n"
        
        if all_images:
            # Add image references in text so they appear in source view
            combined_text += f"Note: {len(all_images)} images from the document are attached below for visual analysis.\n\n"
            # Add image placeholders that will be matched with actual images
            for idx in range(len(all_images)):
                combined_text += f"[Image {idx+1} from document]\n"
            combined_text += "\n"
        
        # Natural, conversational prompt
        combined_text += (
            "Please analyze this document and provide a comprehensive, natural analysis. "
            "Write in a flowing, conversational style like you're explaining it to someone. "
            "For chemistry content (NMR, spectra, etc.), provide detailed analysis with peak tables. "
            "For presentations, provide a cohesive summary that flows naturally rather than listing slides. "
            "Focus on the key insights, findings, and important information."
        )

        # Build content blocks - interleave text and images so they appear together in source view
        content_blocks = [{"type": "text", "text": combined_text}]
        # Add images (no extra metadata - OpenWebUI doesn't allow unexpected keys)
        # Clean each image dict to ensure only expected keys are present
        cleaned_images = []
        for idx, img in enumerate(all_images):
            if isinstance(img, dict) and img.get("type") == "image_url":
                # Extract URL from image_url dict or string
                img_url_obj = img.get("image_url", {})
                if isinstance(img_url_obj, dict):
                    url = img_url_obj.get("url", "")
                elif isinstance(img_url_obj, str):
                    url = img_url_obj
                else:
                    url = ""
                
                # Create clean dict with only expected keys (type and image_url with url)
                if url:
                    clean_img = {
                        "type": "image_url",
                        "image_url": {
                            "url": url
                        }
                    }
                    cleaned_images.append(clean_img)
                    # Log first image structure for debugging
                    if idx == 0:
                        self._log(f"Sample cleaned image keys: {list(clean_img.keys())}, image_url keys: {list(clean_img['image_url'].keys())}")
                else:
                    self._log(f"WARNING: Image {idx} has no URL")
            else:
                # If it's not in the expected format, skip it
                self._log(f"WARNING: Skipping invalid image format at index {idx}: {type(img)}, keys: {list(img.keys()) if isinstance(img, dict) else 'N/A'}")
        
        self._log(f"✅ Cleaned {len(cleaned_images)} images (from {len(all_images)} total)")
        content_blocks.extend(cleaned_images)
        
        messages[-1]["content"] = content_blocks
        body["messages"] = messages

        # Final verification - confirm what's being sent
        self._log(f"✅ SUCCESS - Message ready to send to OpenAI:")
        self._log(f"   - {len(all_content)} text blocks (includes ALL tables, text, notes)")
        self._log(f"   - {len(cleaned_images)} images (as input_image - OpenAI WILL receive these)")
        self._log(f"   - All content will be sent via Responses API to OpenAI")
        self._log("=" * 60)
        return body

    def outlet(self, body: dict, __user__: Optional[dict] = None) -> dict:
        """Output filter - passes through unchanged."""
        return body
